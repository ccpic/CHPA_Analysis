from pyparsing import col
import numpy as np
import pandas as pd
from .data_cls import ChpaAnalyzer
from chart_class import (
    COLOR_DICT,
    PlotLine,
    PlotStackedBar,
    PlotStackedBarPlus,
    PlotHeatGrid,
)
import matplotlib as mpl
import matplotlib.pyplot as plt
import datetime
from matplotlib.gridspec import GridSpec
from pptx import presentation, Presentation, table
from pptx.shapes import picture
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from datetime import datetime
from typing import Tuple, Union

# 预设的位置和宽高参数
IMAGE_TOP = Cm(3.5)
IMAGE_HEIGHT = Cm(13.22)
LABEL_TOP = Cm(2.9)
LABEL_LEFT1 = Cm(25.23)
LABEL_LEFT2 = Cm(28.11)
LABEL_LEFT3 = Cm(30.99)
LABEL_WIDTH = Cm(2.88)
LABEL_HEIGHT = Cm(0.69)

# 预设的body右上角标签
LABELS = [
    {
        "type": "shape",
        # "text": labels[0],
        "top": LABEL_TOP,
        "left": LABEL_LEFT1,
        "width": LABEL_WIDTH,
        "height": LABEL_HEIGHT,
        "bgcolor": RGBColor(0, 128, 128),
    },
    {
        "type": "shape",
        # "text": labels[1],
        "top": LABEL_TOP,
        "left": LABEL_LEFT2,
        "width": LABEL_WIDTH,
        "height": LABEL_HEIGHT,
        "bgcolor": RGBColor(220, 20, 60),
    },
    {
        "type": "shape",
        # "text": labels[2],
        "top": LABEL_TOP,
        "left": LABEL_LEFT3,
        "width": LABEL_WIDTH,
        "height": LABEL_HEIGHT,
        "bgcolor": RGBColor(0, 0, 128),
    },
]


class ChpaPPT(object):
    """Powerpoint幻灯类，包含增加不同内容slide，并按参数添加图片和形状的方法"""

    def __init__(
        self, analyzer: ChpaAnalyzer, template_path: str, save_path: str
    ) -> None:
        """初始化参数
        Parameters
        ----------
        analyzer : MonthlySalesAnalyzer
            一个月度销售分析模块的实例
        template_path : str
            PPT模板路径
        save_path : str
            新PPT保存路径
        """
        self.analyzer = analyzer
        self.template_path = template_path
        self.save_path = save_path
        self.prs = Presentation(template_path)

    def save(self):
        self.prs.save(self.save_path)
        print("PPT has been saved")

    def add_sep_slide(self, title: str = None, layout_style: int = 1, **kwargs):
        """ "添加新间隔页slide， 间隔页slide包含一个居中的大标题文本框，没有其他内容
        Parameters
        ----------
        title : str, optional
            slide标题, by default None
        layout_style : int, optional
            间隔页模板的index, by default 1
        """

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_style])

        title_placeholder = slide.placeholders[0]
        title_placeholder.text = title
        for para in title_placeholder.text_frame.paragraphs:
            font = para.runs[0].font
            if "font_size" in kwargs:
                font.size = Pt(kwargs["font_size"])
            else:
                font.size = Pt(36)

        print("Page%s" % str(int(self.prs.slides.index(slide)) + 1))

    def add_slide(
        self,
        title: str = "",
        layout_style: int = 0,
        imgs: list = None,
        labels: list = None,
        tables: list = None,
        textboxes: list = None,
    ):
        """添加新内容页， 内容页四种类型的shape元素-imgs/labels/tables/textboxes

        Parameters
        ----------
        title : str, optional
            slide标题, by default ""
        layout_style : int, optional
            内容页模板的index, by default 0

        labels : list
            包含要插入矩形含文字标签参数字典的列表, by default None
            参数 :
                top: 标签形状位置-上边距, 默认为Cm(3.5)
                left: 标签形状位置-左边距, 默认为0
                width: 标签形状宽度，默认为Cm(1)
                height: 标签形状宽度，默认为Cm(1)
                text: 标签文本，str
                font_size: 文本大小，str, 默认为Pt(12)
                font_color: 文本颜色，RGBColor, 默认为RGBColor(0, 0, 0)
                bg_color: 形状背景色，RGBColor, 默认为RGBColor(255, 255, 255)
        textboxes: list
            包含要插入文本框参数字典的列表, by default None
        """

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_style])
        shapes = slide.shapes

        # 标题
        title_shape = shapes.title
        title_shape.text = title

        # 插入矩形含文本标签labels
        if labels is not None:
            for label in labels:
                obj_label = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    top=label.get("top", IMAGE_TOP),
                    left=label.get("left", 0),
                    width=label.get("width", Cm(1)),
                    height=label.get("height", Cm(1)),
                )
                obj_label.text = label["text"]
                obj_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                obj_label.text_frame.paragraphs[0].runs[0].font.size = label.get(
                    "font_size", Pt(12)
                )
                obj_label.text_frame.paragraphs[0].runs[0].font.color.rgb = label.get(
                    "font_color", RGBColor(0, 0, 0)
                )
                obj_label.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                fill = obj_label.fill
                fill.solid()  # 填充颜色前必须有此语句
                fill.fore_color.rgb = label.get("bg_color", RGBColor(255, 255, 255))
                obj_label.line.fill.background()  # 去除边框

        # 插入表格tables
        if tables is not None:
            for table in tables:
                table_data = table.get("data").astype(str)

                shape_table = slide.shapes.add_table(
                    rows=table.get("rows", table_data.shape[0]),
                    cols=table.get("cols", table_data.shape[1]),
                    top=table.get("top", IMAGE_TOP),
                    left=table.get("left", 0),
                    width=table.get("width", self.prs.slide_width),
                    height=table.get("height", Cm(10)),
                )
                obj_table = shape_table.table

                # 写入数据文本
                for i, row in enumerate(obj_table.rows):
                    for j, cell in enumerate(row.cells):
                        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # 单元格文本居中（纵向）
                        obj_table.cell(i, j).text = table_data.iloc[i, j]
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER  # 单元格文本居中（横向）
                            paragraph.runs[0].font.size = table.get("fontsize", Pt(11))
                            if table_data.iloc[i, j][0] == "+":  # 正值着绿色
                                paragraph.runs[0].font.color.rgb = RGBColor(0, 176, 80)
                            elif table_data.iloc[i, j][0] == "-":  # 负值着红色
                                paragraph.runs[0].font.color.rgb = RGBColor(255, 0, 0)

                # 表格总样式
                tbl_pr = shape_table._element.graphic.graphicData.tbl
                tbl_pr[0][
                    -1
                ].text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # 微软内置的样式id"NoStyleTableGrid"

                # 单元格颜色
                if "cell_color" in table:
                    for cell_xy, cell_color in table["cell_color"].items():
                        obj_table.cell(cell_xy[0], cell_xy[1]).fill.solid()
                        obj_table.cell(
                            cell_xy[0], cell_xy[1]
                        ).fill.fore_color.rgb = cell_color

                # 字体颜色
                if "font_color" in table:
                    for cell_xy, font_color in table["font_color"].items():
                        for paragraph in obj_table.cell(
                            cell_xy[0], cell_xy[1]
                        ).text_frame.paragraphs:
                            paragraph.runs[0].font.color.rgb = font_color

                # 处理合并单元格的事宜
                if "merge_cells" in table:
                    for merge_cell in table["merge_cells"]:
                        cell1_row = merge_cell[0][0]
                        cell1_col = merge_cell[0][1]
                        cell2_row = merge_cell[1][0]
                        cell2_col = merge_cell[1][1]
                        for i in range(cell2_row, cell1_row - 1, -1):
                            for j in range(cell2_col, cell1_col - 1, -1):
                                if i != cell1_row or j != cell1_col:
                                    obj_table.cell(
                                        i, j
                                    ).text = ""  # 被merge的单元格们只保留左上角主体单元格的文本
                        obj_table.cell(cell1_row, cell1_col).merge(
                            obj_table.cell(cell2_row, cell2_col)
                        )  # 根据指定左上和右下的单元格merge

                if "left" in table is False:
                    shape_table.left = (
                        self.prs.slide_width - shape_table.width
                    ) / 2  # 默认表格居中

        # 插入文本框textboxes
        if textboxes is not None:
            for textbox in textboxes:
                obj_textbox = slide.shapes.add_textbox(
                    top=textbox.get("top", IMAGE_TOP),
                    left=textbox.get("left", 0),
                    width=textbox.get("width", self.prs.slide_width),
                    height=textbox.get("height", Cm(1)),
                )
                text_frame = obj_textbox.text_frame
                p = text_frame.paragraphs[0]
                p.text = textbox.get("text", "")
                p.font.size = textbox.get("font_size", Pt(11))
                p.alignment = textbox.get("alignment", PP_ALIGN.CENTER)

        print("Page%s" % str(int(self.prs.slides.index(slide)) + 1))

    def add_table(
        self,
        data: pd.DataFrame,
        rows: int = None,
        cols: int = None,
        top: Union[float, Inches, Cm] = Cm(3.5),
        left: Union[float, Inches, Cm] = None,
        width: Union[float, Inches, Cm] = None,
        height: Union[float, Inches, Cm] = Cm(10),
        table_style_id: str = "{5940675A-B579-460E-94D1-54222C63F5DA}",
        font_size: Pt = Pt(11),
        cells_color: dict = None,
        fonts_color: dict = None,
        merge_cells: Tuple[Tuple, Tuple] = None,
        col_width: dict = None,
    ) -> table:

        """在当前（最新幻灯页）插入数据表格

        Parameters
        ----------
        data : pd.DataFrame
            表格数据
        rows : int, optional
            表格行数, 如为None怎么默认为data的行数, by default None
        cols : int, optional
            表格列数, 如为None怎么默认为data的列数, by default None
        top : Union[float, Inches, Cm], optional
            表格位置-上边距, by default Cm(3.5)
        left : Union[float, Inches, Cm], optional
            表格位置-左边距, 如为None默认会计算表格自动居中, by default None
        width : Union[float, Inches, Cm], optional
            表格宽度，如为None默认为幻灯片宽度, by default None
        height : Union[float, Inches, Cm], optional
            图片高度, by default Cm(10)
        table_style_id : str, optional
            微软内置表格样式id, 默认为NoStyleTableGrid, by default "{5940675A-B579-460E-94D1-54222C63F5DA}"
        font_size : Pt, optional
            字体大小, by default Pt(11)
        cells_color : dict, optional
            表格特定单元格背景着色, 包含2个参数的tuple, by default None
            参数1: 指定单元格x,y的tuple
            参数2: 指定颜色RGBColor
        fonts_color : dict, optional
            表格特定单元格字体颜色，包含2个参数的tuple, by default None
            参数1: 指定单元格x,y的tuple
            参数2: 指定颜色RGBColor
        merge_cells : Tuple[Tuple, Tuple], optional
            表格特定单元格合并，包含2个参数的tuple, by default None
            参数1: 要合并的最左上的单元格x,y的tuple
            参数2: 要合并的最右下的单元格x,y的tuple
        col_width: dict = None, optional
            指定列的宽度，参数为一个字典：
            key: 指定要调整宽度列的索引
            value: 指定列的宽度
            剩余列宽自动调整
        Returns
        -------
        table
            在幻灯页插入表格后返回表格的形状对象
        """

        slide = self.prs.slides[-1]

        table_data = data.reset_index().T.reset_index().T.astype(str)

        shape_table = slide.shapes.add_table(
            rows=rows if rows is not None else table_data.shape[0],
            cols=cols if cols is not None else table_data.shape[1],
            top=top,
            left=left if left is not None else 0,
            width=width if width is not None else self.prs.slide_width,
            height=height,
        )
        obj_table = shape_table.table

        # 写入数据文本
        for i, row in enumerate(obj_table.rows):
            for j, cell in enumerate(row.cells):
                obj_table.cell(i, j).text = table_data.iloc[i, j]
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER  # 单元格文本居中（横向）
                    paragraph.runs[0].font.size = font_size
                    if table_data.iloc[i, j][0] == "+":  # 正值着绿色
                        paragraph.runs[0].font.color.rgb = RGBColor(0, 176, 80)
                    elif table_data.iloc[i, j][0] == "-":  # 负值着红色
                        paragraph.runs[0].font.color.rgb = RGBColor(255, 0, 0)
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # 单元格文本居中（纵向）
        obj_table.cell(0, 0).text = ""  # 去除因pandas df.reset_index()产生的"index"文本

        # 表格总样式
        tbl_pr = shape_table._element.graphic.graphicData.tbl
        tbl_pr[0][-1].text = table_style_id

        # 单元格颜色
        if cells_color is not None:
            for cell_xy, cell_color in cells_color:
                obj_table.cell(cell_xy[0], cell_xy[1]).fill.solid()
                obj_table.cell(cell_xy[0], cell_xy[1]).fill.fore_color.rgb = cell_color

        # 字体颜色
        if fonts_color is not None:
            for cell_xy, font_color in fonts_color.items():
                for paragraph in obj_table.cell(
                    cell_xy[0], cell_xy[1]
                ).text_frame.paragraphs:
                    paragraph.runs[0].font.color.rgb = font_color

        # 处理合并单元格的事宜
        if merge_cells is not None:
            for merge_cell in merge_cells:
                cell1_row = merge_cell[0][0]
                cell1_col = merge_cell[0][1]
                cell2_row = merge_cell[1][0]
                cell2_col = merge_cell[1][1]
                for i in range(cell2_row, cell1_row - 1, -1):
                    for j in range(cell2_col, cell1_col - 1, -1):
                        if i != cell1_row or j != cell1_col:
                            obj_table.cell(i, j).text = ""  # 被merge的单元格们只保留左上角主体单元格的文本
                obj_table.cell(cell1_row, cell1_col).merge(
                    obj_table.cell(cell2_row, cell2_col)
                )  # 根据指定左上和右下的单元格merge

        # 默认表格居中
        if left is None:
            shape_table.left = int((self.prs.slide_width - shape_table.width) / 2)

        # 设定列宽
        other_width = int(
            (shape_table.width - sum(col_width.values())) / (len(obj_table.columns) -len(col_width))
        )
        for j in range(len(obj_table.columns)):
            if j in col_width:
                obj_table.columns[j].width = col_width[j]
            else:
                obj_table.columns[j].width = other_width

        return shape_table

    def add_image(
        self,
        img_path: str,
        top: Union[float, Inches, Cm] = Cm(3.5),
        left: Union[float, Inches, Cm] = None,
        width: Union[float, Inches, Cm] = None,
        height: Union[float, Inches, Cm] = None,
    ) -> picture:
        """在当前（最新幻灯页）插入图片

        Parameters
        ----------
        img_path : str
            图片路径位置
        top : Union[float, Inches, Cm], optional
            图片位置-上边距, by default Cm(3.5)
        left : Union[float, Inches, Cm], optional
            图片位置-左边距, 如为None默认会计算表格自动居中, by default None
        width : Union[float, Inches, Cm], optional
            图片宽度，如为None且height也为None的情况下会默认为幻灯片宽度, by default None
        height : Union[float, Inches, Cm], optional
            图片高度, by default None

        Returns
        -------
        picture
            在幻灯页插入图片后返回图片的形状对象
        """

        slide = self.prs.slides[-1]

        obj_img = slide.shapes.add_picture(
            image_file=img_path,
            top=top,
            left=left if left is not None else 0,
            width=width
            if width is not None or height is not None
            else self.prs.slide_width,
            height=height,
        )
        if left is None and height is not None:
            obj_img.left = int((self.prs.slide_width - obj_img.width) / 2)  # 默认图片居中

        return obj_img
