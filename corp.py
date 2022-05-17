from re import template
from numpy import product
import pandas as pd
from sqlalchemy import create_engine
from cls.data_cls import ChpaAnalyzer
from cls.ppt_cls import ChpaPPT
from pptx.util import Inches, Pt, Cm
import matplotlib.pyplot as plt
from cls.chart_cls import PlotBubble, PlotWaterfall, PlotStackedBarPlus, PlotStackedBar

if __name__ == "__main__":
    engine = create_engine("mssql+pymssql://(local)/CHPA_1806")
    table_name = "data"

    condition = "([TC I] = 'C CARDIOVASCULAR SYSTEM|心血管系统' \
        and [TC II] not in ('C04 CEREB.+PERIPHE.VASOTHERAP|脑血管和外周血管治疗剂', 'C05 A-VARIC/A-HAEMO PREP|抗静脉曲张/抗痔疮药') \
        and MOLECULE not in ('谷红|ACEGLUTAMIDE+CARTHAMUS TINCTORIUS', '银杏蜜环|ARMILLARIA TABESCENS+GINKGO BILOBA', '棓丙酯|PROPYL GALLATE','葛根素|PUERARIN', '葛根素|PUERARIA')) \
        or [TC II] in ('B01 ANTITHROMBOTIC AGENTS|抗血栓药', 'B02 BLOOD COAG SYST OTH PROD|凝血系统及其他药物') \
        or MOLECULE in ('血脂康胶囊(片)|TRADITIONAL CHINESE MEDICINE')"
    if condition is None:
        sql = f"SELECT * FROM {table_name} WHERE UNIT = 'Value' AND PERIOD = 'MAT'"
    else:
        sql = f"SELECT * FROM {table_name}  WHERE UNIT = 'Value' AND PERIOD = 'MAT' AND ({condition})"

    df = pd.read_sql(sql=sql, con=engine)
    df.to_excel("data.xlsx")
    df["AMOUNT"] = df["AMOUNT"] / 1000000  # 换算单位为百万
    df["PRODUCT"] = df["PRODUCT"].str.split("|").str[0]
    df["MOLECULE"] = df["MOLECULE"].str.split("|").str[0]
    df["PRODUCT_MOLECULE"] = df["PRODUCT"] + "(" + df["MOLECULE"] + ")"
    df["CORPORATION"] = df["CORPORATION"].str.split("|").str[0]

    c = ChpaAnalyzer(data=df, name="CV领域")
    p = ChpaPPT(analyzer=c, template_path="template.pptx", save_path="corp.pptx")

    # CV领域定义市场整体趋势
    market_trend = c.get_pivot(index="DATE")
    market_trend.index = market_trend.index.astype(str)
    market_trend.columns = ["滚动年金额(百万元)"]
    market_gr = market_trend.pct_change(periods=4)
    market_gr.columns = ["同比增长率"]
    market_gr.index = market_gr.index.astype(str)

    f = plt.figure(
        width=15,
        height=6,
        FigureClass=PlotStackedBar,
        fmt=["{:,.0f}"],
        savepath=c.save_path,
        data=market_trend,
        fontsize=10,
        data_line=market_gr,
        fmt_line=["{:,.1%}"],
        style={"title": "CV领域定义市场整体趋势", "xlabel_rotation": 90},
    )

    p.add_slide(title="CV领域定义市场整体趋势")
    p.add_image(f.plot(), height=Cm(13))

    # 所有公司气泡图

    corp_kpi = c.get_kpi("CORPORATION", is_formatted=False)
    corp_kpi = corp_kpi[corp_kpi["滚动年金额"] > 0]

    f = plt.figure(
        width=14,
        height=7,
        FigureClass=PlotBubble,
        gs=None,
        fmt=None,
        savepath=c.save_path,
        data=corp_kpi.loc[:, ["滚动年金额", "净增长", "滚动年金额"]],
        fontsize=10,
        style={
            "title": "CV领域所有企业表现 滚动年金额 vs. 净增长",
            "xlabel": "滚动年金额(百万元)",
            "ylabel": "滚动年金额净增长(百万元)",
        },
    )

    p.add_slide(title="CV领域所有企业表现 滚动年金额 vs. 净增长")
    p.add_image(
        f.plot(label_limit=30, x_fmt="{:,.0f}", y_fmt="{:+,.0f}", y_avg=0),
        width=Cm(30),
        top=Cm(4.14),
    )

    # Top规模公司表格页
    corp_kpi = c.get_kpi("CORPORATION", is_formatted=True)
    product_number = c.get_pivot(
        index="CORPORATION", values="PRODUCT", aggfunc=lambda x: len(x.unique())
    )  # 添加CV领域产品数量统计
    product_number.columns = ["CV领域产品数量"]
    corp_kpi = pd.concat([corp_kpi, product_number], axis=1)

    p.add_slide(title="CV领域Top1-15企业排名及表现明细 最新滚动年- 金额")
    p.add_table(
        corp_kpi.iloc[:15, :],
        width=Cm(30),
        font_size=Pt(9),
        table_style_id="{9D7B26C5-4107-4FEC-AEDC-1716B250A1EF}",
    )

    p.add_slide(title="CV领域Top16-30企业排名及表现明细 最新滚动年- 金额")
    p.add_table(
        corp_kpi.iloc[15:30, :],
        width=Cm(30),
        font_size=Pt(9),
        table_style_id="{9D7B26C5-4107-4FEC-AEDC-1716B250A1EF}",
    )

    # 循环处理Top30公司的表现
    for i, idx in enumerate(corp_kpi.head(30).index):
        rank = i + 1

        # 添加产品表现趋势页
        df_product = c.get_pivot(
            index="DATE",
            columns="PRODUCT_MOLECULE",
            query_str=f"CORPORATION == '{idx}'",
            sort_values="by_last_row",
            ascending=False,
            col_others=10,
        )
        df_product.index = df_product.index.astype(str)

        f = plt.figure(
            width=14,
            height=7,
            FigureClass=PlotStackedBarPlus,
            gs=None,
            fmt=["{:,.0f}"],
            savepath=c.save_path,
            data=df_product.iloc[[-13, -9, -5, -1], :],
            fontsize=12,
            style={
                "title": f"{idx}产品组合表现趋势 - 滚动年 - 金额",
                "ylabel": "滚动年金额(百万元)",
            },
        )

        p.add_slide(title=f"Top{rank} - {idx} - 产品组合表现趋势 - 滚动年 - 金额")
        p.add_image(f.plot(), height=Cm(13))

        # 添加产品净增长贡献页
        product_kpi = c.get_kpi(
            "PRODUCT_MOLECULE",
            query_str=f"CORPORATION == '{idx}'",
            is_formatted=False,
            col_others=10,
        )

        product_kpi["净增长绝对值"] = product_kpi["净增长"].abs()
        product_kpi.sort_values(by="净增长绝对值", ascending=False, inplace=True)
        title = f"{idx}产品组合最新同比净增长贡献 - 滚动年 - 金额"

        value_pre = (
            c.get_pivot(
                "DATE", query_str=f"CORPORATION == '{idx}' and DATE == '2020-12-01'"
            )
            .sum()
            .values[0]
        )

        f = plt.figure(
            width=14,
            height=7,
            FigureClass=PlotWaterfall,
            gs=None,
            fmt=["{:,.0f}"],
            savepath=c.save_path,
            data=product_kpi.loc[:, ["净增长"]],
            fontsize=12,
            style={
                "title": title,
                "xlabel_rotation": 90,
                "remove_yticks": True,
                "ylabel": "滚动年金额净增长(百万元)",
            },
        )

        p.add_slide(title=f"Top{rank} - {idx} - 产品组合最新同比净增长贡献 - 滚动年 - 金额")
        p.add_image(
            f.plot(pre=[("2020", value_pre)], net_index="2021"),
            height=Cm(13),
        )

    p.save()
